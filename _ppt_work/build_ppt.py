"""Build CodiNeo defense PPT at the project root.

Outline (20 slides) follows THESIS 3 DEFENSE- PPT CONTENT.docx:
  1.  Title
  2.  Statement of the Problem
  3.  Objectives
  4.  Scope and Limitations
  5.  Conceptual Framework
  6.  Methodology — Use Case Diagram (image)
  7.  Methodology — Activity Diagram (placeholder)
  8.  Methodology — Class Diagram (placeholder)
  9.  Algorithm — Pipeline Overview
  10. Algorithm — Time & Space Complexity
  11. Chapter 4 — Respondent Profile
  12. Chapter 4 — Functional Suitability
  13. Chapter 4 — Usability + Performance Efficiency
  14. Chapter 4 — Reliability + Maintainability
  15. Chapter 4 — Code Understanding + Doc Outputs
  16. Chapter 4 — Pattern Detection + KPI
  17. Chapter 4 — Cronbach's Alpha (placeholder)
  18. Conclusion (Chapter 5)
  19. Recommendations (Chapter 6)
  20. Thank You / Q&A

Tables follow the BloodLine Ch4 template format (cohort columns, weighted
mean, interpretation row).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Theme ----------
DARK   = RGBColor(0x14, 0x1C, 0x2A)   # deep navy
ACCENT = RGBColor(0x4F, 0x9C, 0xF5)   # cool blue
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xC8, 0xD0, 0xDC)
TEXT   = RGBColor(0x20, 0x28, 0x36)
SOFT   = RGBColor(0xF4, 0xF6, 0xFA)
RULE   = RGBColor(0x4F, 0x9C, 0xF5)

FONT = "Calibri"
SERIF = "Cambria"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # blank

# ---------- helpers ----------
def add_bg(slide, color=SOFT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg

def add_side_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def add_footer(slide, page):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.5), Inches(0.35))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "CodiNeo — A Documentation Generation System for Design Pattern Learning"
    r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = MUTED
    # page number on right
    tb2 = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(page)
    r2.font.name = FONT; r2.font.size = Pt(10); r2.font.color.rgb = MUTED

def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(12.2), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = DARK
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.name = FONT; r2.font.size = Pt(13); r2.font.color.rgb = ACCENT; r2.font.bold = True
    # accent rule under title
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(1.2), Inches(0.06))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background()

def add_bullets(slide, items, left=Inches(0.8), top=Inches(1.7), width=Inches(11.7), height=Inches(5.0),
                size=Pt(16), bold_first_word=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        # bullet
        bullet = p.add_run(); bullet.text = "●  "
        bullet.font.name = FONT; bullet.font.size = size; bullet.font.color.rgb = ACCENT; bullet.font.bold = True
        # text
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run(); r1.text = head + " "
            r1.font.name = FONT; r1.font.size = size; r1.font.bold = True; r1.font.color.rgb = DARK
            r2 = p.add_run(); r2.text = body
            r2.font.name = FONT; r2.font.size = size; r2.font.color.rgb = TEXT
        else:
            r = p.add_run(); r.text = item
            r.font.name = FONT; r.font.size = size; r.font.color.rgb = TEXT

def add_two_col(slide, left_title, left_items, right_title, right_items):
    # left column
    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.5))
    pl = box_l.text_frame.paragraphs[0]
    rl = pl.add_run(); rl.text = left_title
    rl.font.name = FONT; rl.font.size = Pt(16); rl.font.bold = True; rl.font.color.rgb = ACCENT
    add_bullets(slide, left_items, left=Inches(0.8), top=Inches(2.25), width=Inches(6.0), height=Inches(4.6),
                size=Pt(14))
    # right column
    box_r = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.5))
    pr = box_r.text_frame.paragraphs[0]
    rr = pr.add_run(); rr.text = right_title
    rr.font.name = FONT; rr.font.size = Pt(16); rr.font.bold = True; rr.font.color.rgb = ACCENT
    add_bullets(slide, right_items, left=Inches(7.0), top=Inches(2.25), width=Inches(6.0), height=Inches(4.6),
                size=Pt(14))

def add_table(slide, headers, rows, *, left=Inches(0.8), top=Inches(1.8), width=Inches(11.7), row_h=Inches(0.34)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    table = table_shape.table
    # header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
    # body
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else SOFT
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(v)
            r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = TEXT

def add_caption(slide, text, top=Inches(6.6)):
    tb = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = MUTED


# =============================================================================
# Slide 1 — Title
# =============================================================================
s = prs.slides.add_slide(BLANK)
# full-bleed dark background
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
# accent ribbon
ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SW, Inches(0.06))
ribbon.fill.solid(); ribbon.fill.fore_color.rgb = ACCENT; ribbon.line.fill.background()

tb = s.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "CodiNeo"
r.font.name = FONT; r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "A Documentation Generation System Using"
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = MUTED
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "Hash-Based Virtual Structural Copy Algorithm for Design Pattern Learning in DEVCON Luzon"
r3.font.name = FONT; r3.font.size = Pt(20); r3.font.color.rgb = WHITE; r3.font.italic = True

tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.8))
tf2 = tb2.text_frame; tf2.word_wrap = True
for label, name in [
    ("Researchers", "Balbarosa, J. A.  |  De Leon, M. Z.  |  Santander, J. J."),
    ("Institution", "FEU Institute of Technology — Bachelor of Science in Computer Science"),
    ("Defense", "Thesis 3 Defense"),
    ("Date", "2026"),
]:
    p = tf2.add_paragraph()
    r = p.add_run(); r.text = label + ":  "
    r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = name
    r2.font.name = FONT; r2.font.size = Pt(14); r2.font.color.rgb = WHITE
# erase first empty para
tf2.paragraphs[0]._pPr.set("{http://schemas.openxmlformats.org/drawingml/2006/main}lvl", "0")

add_footer(s, 1)

# =============================================================================
# Slide 2 — Statement of the Problem
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Statement of the Problem", "Chapter 1 — Section 1.3")
add_bullets(s, [
    ("General problem:", "Novice C++ developers in collaborative, learning-oriented environments such as DEVCON Luzon struggle to recognise software design patterns inside real source code, and existing tools either skip the structural evidence or emit ungrounded LLM narration."),
    ("Specific problem 1:", "Static-analysis tools surface raw findings without binding them to the design-pattern concept the learner is studying."),
    ("Specific problem 2:", "AI-driven coding assistants generate fluent prose that is not anchored to the actual structure of the submitted code, leading to documentation drift."),
    ("Specific problem 3:", "Existing pattern-detection research is not packaged as a teachable workflow that a student can submit code to and immediately read back evidence-bound explanations."),
    ("Specific problem 4:", "There is no reproducible, evaluation-ready C++ detector that exposes its verdicts, evidence, complexity, and confusion-matrix metrics for academic study within the Philippine collegiate context."),
])
add_footer(s, 2)

# =============================================================================
# Slide 3 — Objectives
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Objectives of the Study", "Chapter 1 — Section 1.4")
add_bullets(s, [
    ("General:", "To design, implement, and evaluate CodiNeo as a documentation-generation system that uses a Hash-Based Virtual Structural Copy algorithm to support design-pattern learning for novice C++ developers in DEVCON Luzon."),
    ("Specific 1:", "Construct a deterministic C++ microservice that performs lexical scan, virtual structural-copy construction, pattern dispatch, evidence binding, and hash-and-link analysis."),
    ("Specific 2:", "Pair every detected verdict with a documentation-oriented explanation and a pre-templated unit-test scaffold drawn from the pattern catalog."),
    ("Specific 3:", "Provide a learner-to-developer pathway through which DEVCON Luzon interns can study a pattern module and then exercise the studio against their own C++ artefacts."),
    ("Specific 4:", "Evaluate CodiNeo against the ISO/IEC 25010 quality model, against KPI-based technical metrics, and against confusion-matrix accuracy for the supported pattern catalog."),
    ("Specific 5:", "Verify the algorithm's time- and space-complexity claims independently through an algorithm evaluator and a research statistician."),
])
add_footer(s, 3)

# =============================================================================
# Slide 4 — Scope and Limitations
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Scope and Limitations", "Chapter 1 — Section 1.7")
add_two_col(s,
    "In scope",
    [
        ("Class-level C++ analysis.", "Complete class and struct declarations."),
        ("Supported pattern catalog.", "Only patterns shipped with the prototype."),
        ("Two pathways.", "Learner modules + Developer Studio."),
        ("DEVCON Luzon cohort.", "Internship-context evaluation."),
        ("Documentation outputs.", "Anchored to bound structural evidence."),
        ("Templated unit tests.", "GoogleTest scaffolds per pattern."),
    ],
    "Out of scope",
    [
        ("Function- or file-only analysis.", "without class context."),
        ("Patterns outside the catalog.", "no claim of generality."),
        ("Runtime tracing.", "verdict is structural, not behavioural."),
        ("Cross-service integration.", "for user-submitted code."),
        ("End-to-end user-flow testing.", "for a class snippet."),
        ("Languages other than C++.", "Java/Python/C# left for future work."),
    ])
add_footer(s, 4)

# =============================================================================
# Slide 5 — Conceptual Framework
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Conceptual Framework", "Chapter 1 — Section 1.6")
# Three-box IPO ribbon
def ipo_box(x_in, label, items, color=DARK):
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(2.1), Inches(3.9), Inches(4.6))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    # label header
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(2.1), Inches(3.9), Inches(0.55))
    head.fill.solid(); head.fill.fore_color.rgb = color; head.line.fill.background()
    htb = head.text_frame; htb.margin_left = Inches(0.2)
    hp = htb.paragraphs[0]
    hr = hp.add_run(); hr.text = label
    hr.font.name = FONT; hr.font.size = Pt(16); hr.font.bold = True; hr.font.color.rgb = WHITE
    # bullets
    tb = s.shapes.add_textbox(Inches(x_in + 0.15), Inches(2.8), Inches(3.6), Inches(3.7))
    tf = tb.text_frame; tf.word_wrap = True
    for k, it in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        bul = p.add_run(); bul.text = "•  "
        bul.font.color.rgb = color; bul.font.bold = True; bul.font.size = Pt(13)
        r = p.add_run(); r.text = it
        r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = TEXT

ipo_box(0.8, "INPUT", [
    "Pasted C++ snippet, file, or multi-file upload",
    "Pattern catalog (token + structural ladder)",
    "DEVCON Luzon intern profile",
    "Expert evaluator profile",
])
ipo_box(4.85, "PROCESS", [
    "Lexical scan & tokenisation",
    "Virtual structural-copy construction",
    "Pattern dispatch (specificity ladder)",
    "Evidence binding & content hash",
    "Documentation + test-scaffold synthesis",
    "AI explanation grounded on verdicts",
])
ipo_box(8.9, "OUTPUT", [
    "Detected pattern verdicts",
    "Bound code-evidence with highlights",
    "Beginner-friendly documentation",
    "Pre-templated unit tests",
    "KPI + confusion-matrix metrics",
    "Saved analysis runs for research",
])
# arrows
for x_in in (4.75, 8.8):
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_in - 0.15), Inches(4.25), Inches(0.35), Inches(0.4))
    ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT; ar.line.fill.background()
add_footer(s, 5)

# =============================================================================
# Slide 6 — Methodology · Use Case Diagram (image)
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Methodology — Use Case Diagram", "Chapter 3 — Section 3.4.4")
USE_CASE_PNG = r"C:\Users\Drew\Desktop\NeoTerritory\docs\Codebase\Diagrams\CodiNeo_UseCase.png"
import os
if os.path.exists(USE_CASE_PNG):
    pic = s.shapes.add_picture(USE_CASE_PNG, Inches(0.8), Inches(1.7), height=Inches(5.0))
    # caption
    add_caption(s, "Figure 5. CodiNeo Use Case Diagram (Learner, Developer, Admin actors with include/extend relations)", top=Inches(6.75))
else:
    add_bullets(s, ["[USE CASE DIAGRAM PNG NOT FOUND — render docs/Codebase/Diagrams/use_case.puml]"])
add_footer(s, 6)

# =============================================================================
# Slide 7 — Methodology · Activity Diagram (placeholder)
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Methodology — Activity Diagram", "Chapter 3 — Section 3.4.3")
# image placeholder box
ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(7.5), Inches(5.0))
ph.fill.solid(); ph.fill.fore_color.rgb = WHITE
ph.line.color.rgb = MUTED; ph.line.width = Pt(1.5); ph.line.dash_style = 7  # dashed
tb = ph.text_frame; tb.word_wrap = True
p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "[ Insert Activity Diagram here ]"
r.font.name = FONT; r.font.size = Pt(16); r.font.color.rgb = MUTED; r.font.italic = True
# narrative on the right
add_bullets(s, [
    ("Login & route.", "Learner / Developer / Admin enter their pathway."),
    ("Learner side.", "Browse modules → study → cross-over to Studio."),
    ("Developer side.", "Submit C++ → validate → analyse → bind evidence → render."),
    ("Save & feedback.", "Persist run, surveys, KPI."),
    ("Admin side.", "Monitor runs, KPI, confusion matrix, exports."),
], left=Inches(8.5), top=Inches(1.7), width=Inches(4.5), height=Inches(5.0), size=Pt(13))
add_footer(s, 7)

# =============================================================================
# Slide 8 — Methodology · Class Diagram (placeholder)
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Methodology — Class Diagram", "Chapter 3 — Section 3.4.5")
ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(7.5), Inches(5.0))
ph.fill.solid(); ph.fill.fore_color.rgb = WHITE
ph.line.color.rgb = MUTED; ph.line.width = Pt(1.5); ph.line.dash_style = 7
tb = ph.text_frame; tb.word_wrap = True
p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "[ Insert Class Diagram here ]"
r.font.name = FONT; r.font.size = Pt(16); r.font.color.rgb = MUTED; r.font.italic = True
add_bullets(s, [
    ("User / TesterAccount.", "Login, role, session."),
    ("LearningModule.", "Pattern lessons & self-check."),
    ("AnalysisRun.", "Submission, hash, verdicts."),
    ("PatternEvidence.", "Bound to code units & lines."),
    ("DocumentationOutput.", "Grounded on PatternEvidence."),
    ("TestScaffold.", "GoogleTest template instance."),
    ("KPIRecord / ConfusionEntry.", "Research-grade metrics."),
], left=Inches(8.5), top=Inches(1.7), width=Inches(4.5), height=Inches(5.0), size=Pt(13))
add_footer(s, 8)

# =============================================================================
# Slide 9 — Algorithm Pipeline (10 stages)
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Algorithm — Hash-Based Virtual Structural Copy Pipeline", "Chapter 3 — Section 3.4")
stages = [
    "1. Source-code submission",
    "2. Input validation (backend)",
    "3. Lexical scanning (microservice)",
    "4. Virtual structural-copy construction",
    "5. Pattern dispatch (specificity ladder)",
    "6. Evidence binding to code units",
    "7. Hash & link analysis (content-addressed)",
    "8. Documentation + test-target preparation",
    "9. Structured JSON report generation",
    "10. Frontend presentation (cards + evidence + tests)",
]
# render as 2-row grid of pills
for i, stage in enumerate(stages):
    col = i % 5
    row = i // 5
    x = Inches(0.8 + col * 2.45)
    y = Inches(2.0 + row * 2.0)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(1.4))
    pill.fill.solid(); pill.fill.fore_color.rgb = WHITE
    pill.line.color.rgb = ACCENT; pill.line.width = Pt(1.5)
    tb = pill.text_frame; tb.word_wrap = True
    tb.margin_left = Inches(0.12); tb.margin_right = Inches(0.12)
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = stage
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = DARK
# footer narrative
tb = s.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Deterministic detector — verdicts are reproducible. AI explanation is grounded on bound evidence and cannot override the detector."
r.font.name = FONT; r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = TEXT
add_footer(s, 9)

# =============================================================================
# Slide 10 — Algorithm Time & Space Complexity
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Algorithm — Time & Space Complexity", "Chapter 3 — Sections 3.4.1 / 3.4.2")
add_two_col(s,
    "Time complexity",
    [
        ("Lexical scan:", "O(N)"),
        ("Structural copy:", "O(N)"),
        ("Pattern dispatch:", "O(N · P) → O(N) since P is fixed"),
        ("Hash & link:", "O(N) average"),
        ("Report generation:", "O(N) bounded by evidence"),
        ("Overall:", "O(N) under prototype scope"),
    ],
    "Space complexity",
    [
        ("Token storage:", "O(N)"),
        ("Structural copy:", "O(N)"),
        ("Hash / link table:", "O(K) ⊆ O(N)"),
        ("Evidence + docs:", "O(E + D) ⊆ O(N)"),
        ("Overall:", "O(N) under prototype scope"),
        ("Verification:", "Algorithm evaluator + research statistician (3.8)"),
    ])
add_footer(s, 10)

# =============================================================================
# Slide 11 — Chapter 4 · Respondent Profile
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Respondent Profile", "Chapter 4 — Section 4.2")
# Intern cohort table (placeholder values blank)
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Intern Cohort (n = ___)"
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Year Level", "Frequency", "Percentage"],
    [["First year","__","__%"], ["Second year","__","__%"], ["Third year","__","__%"],
     ["Fourth year","__","__%"], ["Others","__","__%"], ["Total","__","100%"]],
    left=Inches(0.8), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.32))
# Professional cohort table
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Professional Cohort (n = ___)"
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Years of Experience", "Frequency", "Percentage"],
    [["5–7 years","__","__%"], ["8–10 years","__","__%"], ["11 years and above","__","__%"],
     ["Total","__","100%"]],
    left=Inches(7.0), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.32))
add_caption(s, "Table 10–12. Respondent profile by year level and years of experience. Values pending consolidated export.")
add_footer(s, 11)

# =============================================================================
# Slide 12 — Functional Suitability (BloodLine-style cohort table)
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Functional Suitability", "Chapter 4 — Section 4.3.1")
add_table(s,
    ["Indicator", "Intern Cohort", "Professional Cohort", "Combined Mean", "Interpretation"],
    [
        ["Accuracy of code analysis", "___", "___", "___", "___"],
        ["Relevance of documentation outputs", "___", "___", "___", "___"],
        ["Alignment of evidence with verdicts", "___", "___", "___", "___"],
        ["Usefulness of unit-test targets", "___", "___", "___", "___"],
        ["Weighted Mean", "___", "___", "___", "___"],
    ],
    top=Inches(2.0), row_h=Inches(0.42))
add_caption(s, "Table 14. Functional Suitability of CodiNeo. Likert 1–5. Verbal interpretation per Section 3.8.3.", top=Inches(5.4))
# below the table — interpretation paragraph
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Interpretation. "
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK
r2 = p.add_run(); r2.text = ("Both cohorts assessed the system as performing its intended functions within the supported pattern "
                             "catalog. Mean values, verbal interpretations, and per-question deviations will be populated from the consolidated questionnaire export.")
r2.font.name = FONT; r2.font.size = Pt(13); r2.font.color.rgb = TEXT
add_footer(s, 12)

# =============================================================================
# Slide 13 — Usability + Performance Efficiency (split)
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Usability & Performance Efficiency", "Chapter 4 — Sections 4.3.2 / 4.3.3")
# Usability table
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Usability"
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Indicator", "Mean", "Interpretation"],
    [["Ease of understanding outputs", "___", "___"],
     ["Clarity of documentation", "___", "___"],
     ["Learning-support effectiveness", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(0.8), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.38))
# Performance Efficiency table
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Performance Efficiency"
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Metric", "Value", "Interpretation"],
    [["Average processing time (ms)", "___", "___"],
     ["Memory usage (MB)", "___", "___"],
     ["Throughput (runs / min)", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(7.0), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.38))
add_caption(s, "Tables 15–16. Likert means for usability; KPI measurements for performance. Values pending.", top=Inches(5.4))
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Interpretation. "
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK
r2 = p.add_run(); r2.text = "Usability is reported from intern-cohort responses; performance metrics are pulled directly from the system's KPI store during alpha-test runs. The two are reported as independent silos and are not combined."
r2.font.name = FONT; r2.font.size = Pt(13); r2.font.color.rgb = TEXT
add_footer(s, 13)

# =============================================================================
# Slide 14 — Reliability + Maintainability
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Reliability & Maintainability", "Chapter 4 — Sections 4.3.4 / 4.3.5")
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Reliability"
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Metric", "Mean", "Interpretation"],
    [["Consistency of results", "___", "___"],
     ["Error handling", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(0.8), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.4))
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Maintainability"
r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Metric", "Mean", "Interpretation"],
    [["Ease of extending system", "___", "___"],
     ["Code modularity", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(7.0), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.4))
add_caption(s, "Tables 17–18. Reliability and Maintainability indicators (expert cohort). Values pending.", top=Inches(4.7))
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Interpretation. "
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK
r2 = p.add_run(); r2.text = ("Reliability is interpreted strictly against result-consistency and error-handling behaviour during the evaluation period. Maintainability is interpreted against the layered microservices architecture documented in Section 3.4.6. Each is read on its own.")
r2.font.name = FONT; r2.font.size = Pt(13); r2.font.color.rgb = TEXT
add_footer(s, 14)

# =============================================================================
# Slide 15 — Code Understanding + Documentation Outputs
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Code Understanding & Documentation Outputs", "Chapter 4 — Sections 4.4 / 4.5")
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "System Effectiveness in Code Understanding"
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Indicator", "Mean", "Interpretation"],
    [["Understanding of code structure", "___", "___"],
     ["Identification of key components", "___", "___"],
     ["Comprehension of program logic", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(0.8), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.4))
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Effectiveness of Documentation-Oriented Outputs"
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Indicator", "Mean", "Interpretation"],
    [["Clarity of generated documentation", "___", "___"],
     ["Usefulness for learning", "___", "___"],
     ["Alignment with code structure", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(7.0), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.4))
add_caption(s, "Tables 19–20. Code-understanding and documentation indicators. Values pending.", top=Inches(5.0))
add_footer(s, 15)

# =============================================================================
# Slide 16 — Pattern Detection + KPI
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Pattern Detection & KPI Metrics", "Chapter 4 — Sections 4.6 / 4.7")
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Effectiveness of Design Pattern Detection"
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Indicator", "Mean", "Interpretation"],
    [["Accuracy of pattern identification", "___", "___"],
     ["Helpfulness in learning patterns", "___", "___"],
     ["Weighted Mean", "___", "___"]],
    left=Inches(0.8), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.4))
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Quantitative KPI Metrics (Confusion Matrix)"
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Metric", "Value"],
    [["Accuracy", "___"],
     ["Precision", "___"],
     ["Recall", "___"],
     ["F1 Score", "___"]],
    left=Inches(7.0), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.4))
add_caption(s, "Tables 21–22. Pattern-detection perception (Likert) and technical KPI (confusion-matrix derived). Values pending.", top=Inches(4.7))
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.05), Inches(11.7), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Interpretation. "
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK
r2 = p.add_run(); r2.text = ("The perception metrics (left) describe how respondents experienced the detector. The technical KPIs (right) describe the detector itself, computed against the labelled ground-truth set held constant for every analysis run. The two are read independently.")
r2.font.name = FONT; r2.font.size = Pt(13); r2.font.color.rgb = TEXT
add_footer(s, 16)

# =============================================================================
# Slide 17 — Cronbach's Alpha + Summary of Findings
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Reliability of the Survey & Summary of Findings", "Chapter 4 — Section 4.8 / Chapter 5 — Section 5.2")
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Cronbach's Alpha (Survey Reliability)"
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Subscale", "α"],
    [["Functional Suitability", "___"],
     ["Usability", "___"],
     ["Performance Efficiency", "___"],
     ["Reliability", "___"],
     ["Maintainability", "___"],
     ["Whole survey", "___"]],
    left=Inches(0.8), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.32))
tb = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(0.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Summary of Findings (ISO/IEC 25010)"
r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT
add_table(s,
    ["Category", "Weighted Mean", "Response"],
    [["Functional Suitability", "___", "___"],
     ["Usability", "___", "___"],
     ["Performance Efficiency", "___", "___"],
     ["Reliability", "___", "___"],
     ["Maintainability", "___", "___"],
     ["Total Average", "___", "___"]],
    left=Inches(7.0), top=Inches(2.15), width=Inches(6.0), row_h=Inches(0.32))
add_caption(s, "Tables 23–24. Cronbach's α per subscale (target ≥ 0.70) and ISO/IEC 25010 summary. Values pending.", top=Inches(4.8))
add_footer(s, 17)

# =============================================================================
# Slide 18 — Conclusion
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Conclusions", "Chapter 5 — Summary of Findings and Conclusions")
add_bullets(s, [
    ("1.", "CodiNeo provides a functionally suitable platform for deterministic detection of supported design-pattern evidence in submitted C++ source code, with documentation anchored to that evidence."),
    ("2.", "The Hash-Based Virtual Structural Copy algorithm operates within linear time and linear space bounds, jointly supported by the algorithm evaluator's signed verification and the statistician's regression."),
    ("3.", "Documentation-oriented outputs produced by the AI explanation layer, grounded on deterministic verdicts, are perceived as usefully aligned with the structural evidence they describe."),
    ("4.", "The learner-to-developer pathway allows conceptual exposure to a pattern module to be followed by structural analysis of the learner's own C++ artefacts within the DEVCON Luzon context."),
    ("5.", "The layered microservices architecture and the modular pattern catalog allow CodiNeo to be extended with additional design patterns without re-architecting the analyser."),
])
add_footer(s, 18)

# =============================================================================
# Slide 19 — Recommendations
# =============================================================================
s = prs.slides.add_slide(BLANK); add_bg(s); add_side_bar(s)
add_title(s, "Recommendations", "Chapter 6 — Recommendations")
add_bullets(s, [
    ("Future researchers.", "Expand the pattern catalog beyond the prototype set; repeat the empirical complexity study; port the algorithm to Java, Python, and C# to test generality."),
    ("Instructors & DEVCON Luzon.", "Integrate the learner pathway into the internship curriculum as a pre-requisite for the design-pattern unit."),
    ("Learners.", "Submit one class or struct at a time when first learning a pattern so the analyser's verdict and bound evidence remain legible."),
    ("Maintainers & contributors.", "Keep the catalog as a specificity ladder; keep hashing over the virtual structural copy; keep the deterministic detector authoritative over the AI layer."),
    ("Adopting institutions.", "Run an independent algorithm-verification cycle (evaluator review + statistician regression) before using CodiNeo as a grading input."),
])
add_footer(s, 19)

# =============================================================================
# Slide 20 — Thank You / Q&A
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
ribbon = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.0), SW, Inches(0.06))
ribbon.fill.solid(); ribbon.fill.fore_color.rgb = ACCENT; ribbon.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank you"
r.font.name = FONT; r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Questions?"
r2.font.name = FONT; r2.font.size = Pt(32); r2.font.color.rgb = ACCENT

OUT = r"C:\Users\Drew\Desktop\NeoTerritory\CodiNeo - Thesis 3 Defense.pptx"
prs.save(OUT)
print("Wrote:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst) if hasattr(prs.slides, '_sldIdLst') else "ok")
